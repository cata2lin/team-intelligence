# -*- coding: utf-8 -*-
"""Helpere reutilizabile pentru prezentări .pptx în stil „template Canva"
(crem/foto-arcuit/numere mari estompate). Testat cu python-pptx 1.0.2 pe macOS.

Trucuri cheie incluse:
  - colț arcuit pe poze (mască PIL cu pieslice → PNG cu alpha)
  - overlay semitransparent pe full-bleed (hack XML <a:alpha> pe srgbClr)
  - letter-spacing pe titluri mono (rPr.set('spc', ...))
  - numere mari „tăiate" de marginea slide-ului
  - shadow.inherit=False (altfel apar umbre nedorite pe forme)
  - export PDF de verificare prin PowerPoint (open -a + AppleScript)

Fonturi: instalează-le întâi în ~/Library/Fonts (Google Fonts), altfel PowerPoint
substituie. Ex.: Montserrat (repo JulietaUla), Anonymous Pro (google/fonts).
"""
import os
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Paletă implicită (suprascrie per proiect)
CREAM = RGBColor(0xEF, 0xE9, 0xE1)
INK   = RGBColor(0x2F, 0x2A, 0x26)
TEAL  = RGBColor(0x33, 0x58, 0x5C)
NUMC  = RGBColor(0xDD, 0xD4, 0xC8)
MUTED = RGBColor(0x7A, 0x70, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MONO, SANS, SANS_SB = "Anonymous Pro", "Montserrat", "Montserrat SemiBold"
SW, SH = Inches(13.333), Inches(7.5)


# ---------- pregătire imagini (PIL) ----------

def cover_169(src, out, W=1700):
    """Crop centrat la 16:9 pentru fundaluri full-bleed."""
    im = Image.open(src).convert("RGB")
    w, h = im.size
    tw = int(h * 16 / 9)
    if tw <= w:
        x0 = (w - tw) // 2; im = im.crop((x0, 0, x0 + tw, h))
    else:
        th = int(w * 9 / 16); y0 = (h - th) // 2; im = im.crop((0, y0, w, y0 + th))
    im.resize((W, int(W * 9 / 16)), Image.LANCZOS).save(out, quality=87)
    return out


def arch_photo(src, out, corner="tr", ratio=0.765, H=1500, radius_frac=0.285):
    """Poză verticală cu UN colț arcuit mare (ca în template-urile Canva).
    corner: 'tr' (poză pe stânga slide-ului) sau 'tl' (poză pe dreapta)."""
    im = Image.open(src).convert("RGB")
    w, h = im.size
    tw = int(h * ratio)
    if tw <= w:
        x0 = (w - tw) // 2; im = im.crop((x0, 0, x0 + tw, h))
    else:
        th = int(w / ratio); y0 = (h - th) // 2; im = im.crop((0, y0, w, y0 + th))
    W = int(H * ratio)
    im = im.resize((W, H), Image.LANCZOS)
    R = int(H * radius_frac)
    mask = Image.new("L", (W, H), 255)
    d = ImageDraw.Draw(mask)
    if corner == "tr":
        d.rectangle([W - R, 0, W, R], fill=0)
        d.pieslice([W - 2 * R, 0, W, 2 * R], 270, 360, fill=255)
    else:
        d.rectangle([0, 0, R, R], fill=0)
        d.pieslice([0, 0, 2 * R, 2 * R], 180, 270, fill=255)
    rgba = im.convert("RGBA"); rgba.putalpha(mask); rgba.save(out)
    return out


# ---------- primitive pptx ----------

def deck():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    return prs


def slide_new(prs, bg=CREAM):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def txt(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tb, tf


def run(p, text, font=SANS, size=16, bold=False, color=INK, italic=False, spc=None):
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if spc: r.font._rPr.set('spc', str(spc))  # letter-spacing, sutimi de pt
    return r


def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    if line: p.line_spacing = line
    return p


def full_photo(s, img, alpha=45):
    """Fundal full-bleed + overlay întunecat semitransparent (hack a:alpha)."""
    s.shapes.add_picture(img, 0, 0, SW, SH)
    ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0x14, 0x0E, 0x0A)
    ov.line.fill.background(); ov.shadow.inherit = False
    srgb = ov.fill.fore_color._xFill.find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)}))


def white_card(s, x, y, w, h, radius=0.055):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = radius
    c.fill.solid(); c.fill.fore_color.rgb = WHITE
    c.line.fill.background(); c.shadow.inherit = False
    return c


def big_num(s, n, x=None, color=NUMC, size=170):
    """Număr mare estompat, parțial sub marginea slide-ului (PowerPoint îl taie)."""
    if x is None: x = SW - Inches(3.1)
    tb, tf = txt(s, x, SH - Inches(2.35), Inches(3.0), Inches(2.6), align=PP_ALIGN.RIGHT)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT, space_after=0, line=1.0)
    run(p, n, font=MONO, size=size, bold=True, color=color)


def kicker_title(s, x, y, w, kicker, title, tsize=30):
    tb, tf = txt(s, x, y, w, Inches(1.5))
    p = para(tf, first=True, space_after=4)
    run(p, kicker.upper(), font=SANS_SB, size=13, bold=True, color=MUTED, spc=120)
    p2 = para(tf, space_after=0, line=1.0)
    run(p2, title.upper(), font=MONO, size=tsize, bold=True, color=INK, spc=200)


def bullets(s, x, y, w, h, items, size=15.5, gap=10, line=1.12, accent=TEAL):
    """items: str sau (lead_bold, rest)."""
    tb, tf = txt(s, x, y, w, h)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap, line=line)
        run(p, "▪  ", font=SANS, size=size, bold=True, color=accent)
        if isinstance(it, tuple):
            run(p, it[0], font=SANS_SB, size=size, bold=True, color=INK)
            run(p, it[1], font=SANS, size=size, color=INK)
        else:
            run(p, it, font=SANS, size=size, color=INK)


def table_card(s, x, y, w, h, header, rows, widths=None, head_fill=TEAL,
               head_size=12.5, cell_size=11.5, radius=0.04):
    """Tabel REAL (editabil) pe un card alb — pt calendare, liste de articole,
    grile de cuvinte-cheie. widths = fracții care însumează 1 (altfel egale)."""
    white_card(s, x, y, w, h, radius=radius)
    from pptx.util import Emu
    iw = w - Inches(0.4)
    tbl = s.shapes.add_table(len(rows) + 1, len(header),
                             x + Inches(0.2), y + Inches(0.2), iw, h - Inches(0.4)).table
    tbl.first_row = False
    widths = widths or [1.0 / len(header)] * len(header)
    for j, fr in enumerate(widths):
        tbl.columns[j].width = Emu(int(iw * fr))

    def cell(r, c, text, bold=False, color=INK, size=cell_size, fill=WHITE):
        cl = tbl.cell(r, c)
        cl.margin_left = cl.margin_right = Inches(0.08)
        cl.margin_top = cl.margin_bottom = Inches(0.04)
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
        p = cl.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run(p, text, font=SANS_SB if bold else SANS, size=size, bold=bold, color=color)
    for j, htext in enumerate(header):
        cell(0, j, htext, bold=True, color=WHITE, size=head_size, fill=head_fill)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell(i + 1, j, str(v))
    return tbl


def export_pdf_via_powerpoint(pptx_path, pdf_path):
    """Preview fidel: PowerPoint refuză uneori `open` din AppleScript (−9074);
    deschide întâi cu `open -a`, apoi salvează prin AppleScript."""
    os.system(f'pkill -9 -f "Microsoft PowerPoint" 2>/dev/null; sleep 2')
    os.system(f'open -a "Microsoft PowerPoint" "{pptx_path}"; sleep 12')
    script = (f'with timeout of 240 seconds\n tell application "Microsoft PowerPoint"\n'
              f'  save active presentation in (POSIX file "{pdf_path}") as save as PDF\n'
              f' end tell\nend timeout')
    os.system(f"osascript -e '{script}'")
    os.system('pkill -9 -f "Microsoft PowerPoint" 2>/dev/null')
    return os.path.exists(pdf_path)
